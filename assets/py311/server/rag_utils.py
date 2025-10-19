
# import required modules
import os
import time
import docx
import olefile
#import hwp5
import PyPDF2
import logging
from logging.handlers import TimedRotatingFileHandler

# Extract text from pptx files using python-pptx
from pptx import Presentation
# Read HWPX file (XML-based Hangul document)
import xml.etree.ElementTree as ET

DEBUG_MODE = False
BEGIN_TIME = time.time()

_logger = None
FILE_LOGGING_ENABLED = False

def get_logger():
    global _logger

    if _logger is not None:
        return _logger

    # Logger setup
    logger = logging.getLogger("CogApp")
    logger.setLevel(logging.DEBUG)

    _logger = logger

    return logger


def is_debug_mode() -> bool:
    global DEBUG_MODE
    return DEBUG_MODE


def enable_file_logging(log_file_path: str, log_level: int = logging.DEBUG):
    global FILE_LOGGING_ENABLED
    global DEBUG_MODE

    if FILE_LOGGING_ENABLED:
        return  # 이미 파일 로깅이 활성화된 경우

    log_file_path = log_file_path.replace("/", os.sep)
    DEBUG_MODE = True

    # if log_file_path is not writable:
    if not os.access(os.path.dirname(log_file_path), os.W_OK):
        # create directory if not exists
        os.makedirs(os.path.dirname(log_file_path), exist_ok=True)

    FILE_LOGGING_ENABLED = True
    logger = get_logger()
    # TimedRotatingFileHandler로 교체 (매일 자정마다 rotate)
    handler = TimedRotatingFileHandler(log_file_path, when="midnight", interval=1, backupCount=7, encoding="utf-8")
    handler.suffix = "%Y-%m-%d"  # 파일명에 날짜 추가
    handler.setLevel(log_level)
    formatter = logging.Formatter("%(asctime)s - %(name)s - %(levelname)s - %(message)s")
    handler.setFormatter(formatter)
    logger.addHandler(handler)


def print_elapsed_time_and_update(tag: str, begin_time: float = None):
    global BEGIN_TIME
    
    logger = get_logger()

    if begin_time is not None:
        BEGIN_TIME = begin_time
        #print(f"BEGIN_TIME initialized to {BEGIN_TIME}")
        return
    assert BEGIN_TIME > 0, "BEGIN_TIME이 초기화되지 않았습니다."

    if is_debug_mode():
        logger.debug(f"Time for {tag}: {time.time() - BEGIN_TIME:.2f} seconds")
    BEGIN_TIME = time.time()


def read_hwpx_file(file_path):
    """
    Extracts all text from a .hwpx file (HWPX XML format).
    Args:
        file_path (str): Path to the .hwpx file.
    Returns:
        str: Extracted text content.
    Raises:
        ValueError: If the file is not a .hwpx file.
        Exception: For other errors during extraction.
    """
    if not file_path.lower().endswith('.hwpx'):
        raise ValueError("Only .hwpx files are supported.")
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        ns = {'h': 'http://www.hancom.co.kr/hwpml/2010/10/20'}
        paragraphs = root.findall('.//h:p', ns)
        text = []
        for para in paragraphs:
            runs = para.findall('.//h:t', ns)
            for run in runs:
                if run.text:
                    text.append(run.text)
        return '\n'.join(text)
    except Exception as e:
        raise Exception(f"Failed to extract text from hwpx: {e}")


def extract_text_from_pptx(file_path):
    """
    Extracts all text from a .pptx file.
    Args:
        file_path (str): Path to the .pptx file.
    Returns:
        str: Extracted text content.
    Raises:
        ValueError: If the file is not a .pptx file.
        Exception: For other errors during extraction.
    """
    if not file_path.lower().endswith('.pptx'):
        raise ValueError("Only .pptx files are supported.")
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                    text_runs.append(shape.text_frame.text)
        return "\n".join([t for t in text_runs if t.strip()])
    except Exception as e:
        raise Exception(f"Failed to extract text from pptx: {e}")


# Utility function to read plain text files
def read_text_file(file_path, encoding="utf-8"):
    """
    Reads the content of a plain text file.
    Args:
        file_path (str): Path to the text file.
        encoding (str): File encoding (default: utf-8).
    Returns:
        str: File content as a string.
    Raises:
        FileNotFoundError: If the file does not exist.
        Exception: For other errors during reading.
    """
    try:
        with open(file_path, "r", encoding=encoding) as f:
            return f.read()
    except FileNotFoundError:
        raise
    except Exception as e:
        raise Exception(f"Failed to read text file: {e}")


def read_msword_file(file_path):

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)


def read_pdf_file(file_path):

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    full_text = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            full_text.append(page.extract_text())
    return '\n'.join(full_text)


"""
def read_hwp_file(file_path):

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    with olefile.OleFileIO(file_path) as ole:
        if not ole.exists('PrvText'):
            raise ValueError("HWP 파일에 PrvText 스트림이 없습니다.")
        stream = ole.openstream('PrvText')
        hwp = hwp5.HWP5File(stream)
        text = []
        for para in hwp.bodytext:
            text.append(para.text)
    return '\n'.join(text)
"""


def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.txt':
        return read_text_file(file_path)
    elif ext == '.docx':
        return read_msword_file(file_path)
    elif ext == '.pdf':
        return read_pdf_file(file_path)
    #elif ext == '.hwp':
    #    return read_hwp_file(file_path)
    elif ext == '.hwpx':
        return read_hwpx_file(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"지원되지 않는 파일 형식입니다: {ext}")


def write_text_to_file(file_path, text):
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text)
    #print(f"텍스트가 {file_path}에 저장되었습니다.")


# ===== 유틸: 파일 로드 & 청크 =====
def read_text(path):
    """
    읽을 파일을 열어 텍스트를 반환합니다.

    Parameters
    ----------
    path : str
        파일 경로

    Returns
    -------
    str
        읽은 텍스트

    Raises
    ------
    FileNotFoundError
        파일이 존재하지 않을 때
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"지식 파일을 찾을 수 없습니다: {path}")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

